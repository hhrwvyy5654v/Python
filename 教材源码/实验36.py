import pptx

# 打开已有演示文档
pptFile = pptx.Presentation('test.pptx')
# 遍历所有幻灯片
for slide in pptFile.slides:
    # 遍历当前幻灯片中所有组件
    for shape in slide.shapes:
        # 找到表格
        if shape.shape_type==19:
            table = shape
            # 遍历并输出单元格内容
            for row in table.table.rows:
                for cell in row.cells:
                    print(cell.text_frame.text, end='\t')
                print()
            print()
